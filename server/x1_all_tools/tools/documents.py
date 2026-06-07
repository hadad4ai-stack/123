from __future__ import annotations
from pathlib import Path
from typing import Any
import csv, json
from x1_all_tools.registry import ToolSpec, object_schema
from x1_all_tools.security import safe_join

def _minimal_pdf(path: Path, title: str, paragraphs: list[str]) -> None:
    text = title + "\n\n" + "\n\n".join(str(p) for p in paragraphs)
    # Minimal single-page PDF; ASCII-safe fallback.
    safe = text.encode("latin-1", errors="replace").decode("latin-1")
    lines = safe.splitlines()[:40]
    stream_lines = ["BT", "/F1 14 Tf", "72 760 Td"]
    for i, line in enumerate(lines):
        escaped = line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        if i:
            stream_lines.append("0 -18 Td")
        stream_lines.append(f"({escaped[:90]}) Tj")
    stream_lines.append("ET")
    stream = "\n".join(stream_lines)
    objects = [
        "1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj",
        "2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj",
        "3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >> endobj",
        "4 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj",
        f"5 0 obj << /Length {len(stream.encode('latin-1'))} >> stream\n{stream}\nendstream endobj",
    ]
    pdf = "%PDF-1.4\n"
    offsets = []
    for obj in objects:
        offsets.append(len(pdf.encode("latin-1")))
        pdf += obj + "\n"
    xref = len(pdf.encode("latin-1"))
    pdf += f"xref\n0 {len(objects)+1}\n0000000000 65535 f \n"
    for off in offsets:
        pdf += f"{off:010d} 00000 n \n"
    pdf += f"trailer << /Size {len(objects)+1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF"
    path.write_bytes(pdf.encode("latin-1"))

def pdf_create(title: str, paragraphs: list[str], path: str = "report.pdf", runtime=None) -> dict[str, Any]:
    target = safe_join(runtime.workspace, path)
    target.parent.mkdir(parents=True, exist_ok=True)
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import cm
        doc = SimpleDocTemplate(str(target), pagesize=A4, rightMargin=2*cm, leftMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)
        styles = getSampleStyleSheet()
        story = [Paragraph(title, styles["Title"]), Spacer(1, 0.5*cm)]
        for p in paragraphs:
            story.append(Paragraph(str(p).replace("\n", "<br/>"), styles["BodyText"]))
            story.append(Spacer(1, 0.25*cm))
        doc.build(story)
        engine = "reportlab"
    except Exception:
        _minimal_pdf(target, title, paragraphs)
        engine = "minimal_pdf_fallback"
    return {"path": str(target), "bytes": target.stat().st_size, "engine": engine}

def _pypdf():
    try:
        import pypdf
        return pypdf
    except ImportError as exc:
        raise RuntimeError("This PDF tool requires pypdf: pip install pypdf") from exc

def pdf_extract_text(path: str, max_chars: int = 100000, runtime=None) -> dict[str, Any]:
    pypdf = _pypdf()
    target = safe_join(runtime.workspace, path)
    reader = pypdf.PdfReader(str(target))
    pages = []
    total = ""
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        pages.append({"page": i + 1, "text": text})
        total += text + "\n"
        if len(total) >= max_chars:
            break
    return {"path": str(target), "pages": pages, "text": total[:max_chars], "truncated": len(total) > max_chars}

def pdf_merge(paths: list[str], output: str, runtime=None) -> dict[str, Any]:
    pypdf = _pypdf()
    out = safe_join(runtime.workspace, output)
    out.parent.mkdir(parents=True, exist_ok=True)
    writer = pypdf.PdfWriter()
    count = 0
    for p in paths:
        reader = pypdf.PdfReader(str(safe_join(runtime.workspace, p)))
        for page in reader.pages:
            writer.add_page(page)
            count += 1
    with out.open("wb") as f:
        writer.write(f)
    return {"path": str(out), "pages": count, "bytes": out.stat().st_size}

def pdf_split(path: str, output_dir: str = "pdf_pages", runtime=None) -> dict[str, Any]:
    pypdf = _pypdf()
    src = safe_join(runtime.workspace, path)
    outdir = safe_join(runtime.workspace, output_dir)
    outdir.mkdir(parents=True, exist_ok=True)
    reader = pypdf.PdfReader(str(src))
    outputs = []
    for i, page in enumerate(reader.pages):
        writer = pypdf.PdfWriter()
        writer.add_page(page)
        out = outdir / f"page_{i+1}.pdf"
        with out.open("wb") as f:
            writer.write(f)
        outputs.append(str(out.relative_to(runtime.workspace)))
    return {"source": str(src), "pages": len(outputs), "outputs": outputs}

def pdf_compress(path: str, output: str, runtime=None) -> dict[str, Any]:
    pypdf = _pypdf()
    src = safe_join(runtime.workspace, path)
    out = safe_join(runtime.workspace, output)
    out.parent.mkdir(parents=True, exist_ok=True)
    reader = pypdf.PdfReader(str(src))
    writer = pypdf.PdfWriter()
    for page in reader.pages:
        try:
            page.compress_content_streams()
        except Exception:
            pass
        writer.add_page(page)
    with out.open("wb") as f:
        writer.write(f)
    return {"path": str(out), "original_bytes": src.stat().st_size, "bytes": out.stat().st_size}

def pdf_images(path: str, output_dir: str = "pdf_images", runtime=None) -> dict[str, Any]:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("pdf.images requires PyMuPDF: pip install pymupdf") from exc
    src = safe_join(runtime.workspace, path)
    outdir = safe_join(runtime.workspace, output_dir)
    outdir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(src))
    outputs = []
    for page_index in range(len(doc)):
        for img_index, img in enumerate(doc[page_index].get_images(full=True)):
            xref = img[0]
            base = doc.extract_image(xref)
            ext = base.get("ext", "png")
            out = outdir / f"page_{page_index+1}_img_{img_index+1}.{ext}"
            out.write_bytes(base["image"])
            outputs.append(str(out.relative_to(runtime.workspace)))
    return {"source": str(src), "images": outputs}

def docx_create(title: str, paragraphs: list[str], path: str = "document.docx", runtime=None) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("docx.create requires python-docx: pip install python-docx") from exc
    target = safe_join(runtime.workspace, path)
    target.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    doc.add_heading(title, 0)
    for p in paragraphs:
        doc.add_paragraph(str(p))
    doc.save(target)
    return {"path": str(target), "bytes": target.stat().st_size}

def docx_read(path: str, runtime=None) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("docx.read requires python-docx: pip install python-docx") from exc
    target = safe_join(runtime.workspace, path)
    doc = Document(str(target))
    paragraphs = [p.text for p in doc.paragraphs]
    return {"path": str(target), "paragraphs": paragraphs, "text": "\n".join(paragraphs)}

def docx_edit(path: str, output: str, find: str, replace: str, runtime=None) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("docx.edit requires python-docx: pip install python-docx") from exc
    src = safe_join(runtime.workspace, path)
    out = safe_join(runtime.workspace, output)
    out.parent.mkdir(parents=True, exist_ok=True)
    doc = Document(str(src))
    replacements = 0
    for p in doc.paragraphs:
        if find in p.text:
            p.text = p.text.replace(find, replace)
            replacements += 1
    doc.save(out)
    return {"path": str(out), "replacements": replacements, "bytes": out.stat().st_size}

def xlsx_create(path: str, sheets: list[dict[str, Any]], runtime=None) -> dict[str, Any]:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
        from openpyxl.utils import get_column_letter
    except ImportError as exc:
        raise RuntimeError("xlsx.create requires openpyxl: pip install openpyxl") from exc
    target = safe_join(runtime.workspace, path)
    target.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    wb.remove(wb.active)
    for sheet in sheets:
        ws = wb.create_sheet(str(sheet.get("name", "Sheet"))[:31])
        rows = sheet.get("rows", [])
        for row in rows:
            ws.append(list(row))
        if rows:
            for cell in ws[1]:
                cell.font = Font(bold=True)
                cell.fill = PatternFill("solid", fgColor="DDDDDD")
            for col in ws.columns:
                width = min(max(len(str(c.value)) if c.value is not None else 0 for c in col) + 2, 60)
                ws.column_dimensions[get_column_letter(col[0].column)].width = width
    if not wb.sheetnames:
        wb.create_sheet("Sheet1")
    wb.save(target)
    return {"path": str(target), "sheets": wb.sheetnames, "bytes": target.stat().st_size}

def xlsx_read(path: str, sheet: str | None = None, runtime=None) -> dict[str, Any]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("xlsx.read requires openpyxl: pip install openpyxl") from exc
    target = safe_join(runtime.workspace, path)
    wb = load_workbook(target, data_only=True)
    ws = wb[sheet] if sheet else wb[wb.sheetnames[0]]
    rows = [[cell.value for cell in row] for row in ws.iter_rows()]
    return {"path": str(target), "sheet": ws.title, "rows": rows}

def xlsx_write(path: str, sheet: str, cell: str, value: Any, output: str | None = None, runtime=None) -> dict[str, Any]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("xlsx.write requires openpyxl: pip install openpyxl") from exc
    target = safe_join(runtime.workspace, path)
    out = safe_join(runtime.workspace, output or path)
    wb = load_workbook(target)
    ws = wb[sheet] if sheet in wb.sheetnames else wb.create_sheet(sheet)
    ws[cell] = value
    wb.save(out)
    return {"path": str(out), "sheet": sheet, "cell": cell, "value": value}

def pptx_create(title: str, slides: list[dict[str, Any]], path: str = "presentation.pptx", runtime=None) -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError as exc:
        raise RuntimeError("pptx.create requires python-pptx: pip install python-pptx") from exc
    target = safe_join(runtime.workspace, path)
    target.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated by X1 All Tools"
    for item in slides:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = str(item.get("title", "Slide"))
        tf = s.placeholders[1].text_frame
        tf.clear()
        bullets = item.get("bullets", [])
        body = item.get("body")
        if body:
            p = tf.paragraphs[0]
            p.text = str(body)
            p.font.size = Pt(20)
        for b in bullets:
            p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
            p.text = str(b)
            p.font.size = Pt(20)
    prs.save(target)
    return {"path": str(target), "slides": len(prs.slides), "bytes": target.stat().st_size}

def pptx_edit(path: str, output: str, find: str, replace: str, runtime=None) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("pptx.edit requires python-pptx: pip install python-pptx") from exc
    src = safe_join(runtime.workspace, path)
    out = safe_join(runtime.workspace, output)
    prs = Presentation(str(src))
    replacements = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and find in shape.text:
                shape.text = shape.text.replace(find, replace)
                replacements += 1
    prs.save(out)
    return {"path": str(out), "replacements": replacements, "bytes": out.stat().st_size}

TOOLS = [
    ToolSpec("pdf.create", "Create a PDF document.", object_schema({"title": {"type": "string"}, "paragraphs": {"type": "array"}, "path": {"type": "string", "default": "report.pdf"}}, ["title", "paragraphs"]), pdf_create),
    ToolSpec("pdf.extract_text", "Extract text from PDF.", object_schema({"path": {"type": "string"}, "max_chars": {"type": "integer", "default": 100000}}, ["path"]), pdf_extract_text),
    ToolSpec("pdf.split", "Split PDF into single-page PDFs.", object_schema({"path": {"type": "string"}, "output_dir": {"type": "string", "default": "pdf_pages"}}, ["path"]), pdf_split),
    ToolSpec("pdf.merge", "Merge multiple PDF files.", object_schema({"paths": {"type": "array"}, "output": {"type": "string"}}, ["paths", "output"]), pdf_merge),
    ToolSpec("pdf.compress", "Compress PDF content streams.", object_schema({"path": {"type": "string"}, "output": {"type": "string"}}, ["path", "output"]), pdf_compress),
    ToolSpec("pdf.images", "Extract images from a PDF using PyMuPDF.", object_schema({"path": {"type": "string"}, "output_dir": {"type": "string", "default": "pdf_images"}}, ["path"]), pdf_images),
    ToolSpec("docx.create", "Create Word DOCX.", object_schema({"title": {"type": "string"}, "paragraphs": {"type": "array"}, "path": {"type": "string", "default": "document.docx"}}, ["title", "paragraphs"]), docx_create),
    ToolSpec("docx.read", "Read Word DOCX.", object_schema({"path": {"type": "string"}}, ["path"]), docx_read),
    ToolSpec("docx.edit", "Edit DOCX by find/replace.", object_schema({"path": {"type": "string"}, "output": {"type": "string"}, "find": {"type": "string"}, "replace": {"type": "string"}}, ["path", "output", "find", "replace"]), docx_edit),
    ToolSpec("xlsx.create", "Create Excel XLSX.", object_schema({"path": {"type": "string"}, "sheets": {"type": "array"}}, ["path", "sheets"]), xlsx_create),
    ToolSpec("xlsx.read", "Read Excel XLSX.", object_schema({"path": {"type": "string"}, "sheet": {"type": ["string", "null"], "default": None}}, ["path"]), xlsx_read),
    ToolSpec("xlsx.write", "Write a value to Excel XLSX.", object_schema({"path": {"type": "string"}, "sheet": {"type": "string"}, "cell": {"type": "string"}, "value": {}, "output": {"type": ["string", "null"], "default": None}}, ["path", "sheet", "cell", "value"]), xlsx_write),
    ToolSpec("pptx.create", "Create PowerPoint PPTX.", object_schema({"title": {"type": "string"}, "slides": {"type": "array"}, "path": {"type": "string", "default": "presentation.pptx"}}, ["title", "slides"]), pptx_create),
    ToolSpec("pptx.edit", "Edit PowerPoint by find/replace.", object_schema({"path": {"type": "string"}, "output": {"type": "string"}, "find": {"type": "string"}, "replace": {"type": "string"}}, ["path", "output", "find", "replace"]), pptx_edit),
]
